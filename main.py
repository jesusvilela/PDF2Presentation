import nltk
import PyPDF2
import fitz
import openai
import torch
import re
import spacy
from pptx import Presentation
from pptx.util import Inches
from transformers import T5Tokenizer, T5ForConditionalGeneration
from collections import Counter
from io import BytesIO
from PIL import Image
import os
from diffusers import StableDiffusionPipeline, EulerDiscreteScheduler
import time

nltk.download('punkt')
openai.api_key = "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"#"your_openai_key
nlp = spacy.load('en_core_web_sm') #python -m spacy download en_core_web_sm
#record start time
start_time = time.time()

def preprocess_text(text):
    doc = nlp(text)
    cleaned_text = " ".join(token.lemma_ for token in doc if not token.is_stop)
    return cleaned_text

def extract_text_from_pdf(file_path):
    with open(file_path, "rb") as file:
        reader = PyPDF2.PdfFileReader(file)
        text = []
        for page in range(reader.getNumPages()):
            text.append(reader.getPage(page).extract_text())
    return text

def extract_images_from_pdf(file_path):
    doc = fitz.open(file_path)
    images = []
    image_indices = []
    for i in range(len(doc)):
        for img in doc.get_page_images(i):
            xref = img[0]
            img_data = doc.extract_image(xref)
            image = Image.open(BytesIO(img_data['image']))
            # Convert image to PNG
            png_image = BytesIO()
            image.save(png_image, format='PNG')
            images.append(png_image.getvalue())
            image_indices.append(i)
    return images, image_indices

def generate_summary(section):
    model = T5ForConditionalGeneration.from_pretrained('t5-base')
    tokenizer = T5Tokenizer.from_pretrained('t5-base')#pip install sentencepiece
    inputs = tokenizer.encode("summarize: " + section, return_tensors="pt", max_length=512)
    outputs = model.generate(inputs, max_length=150, min_length=40, length_penalty=2.0, num_beams=4, early_stopping=True)
    return tokenizer.decode(outputs[0])

def extract_key_phrases(text):
    doc = nlp(text)
    return [chunk.text for chunk in doc.noun_chunks]

def generate_title(section, summary):
    key_phrases = extract_key_phrases(summary)
    title = ' '.join(key_phrases[:3])  # Use the first 3 key phrases to generate the title
    return title
def generate_cover(title):
    model_id = "stabilityai/stable-diffusion-2-1-base"
    scheduler = EulerDiscreteScheduler.from_pretrained(model_id, subfolder="scheduler")
    pipe = StableDiffusionPipeline.from_pretrained(model_id, scheduler=scheduler, torch_dtype=torch.float32)
    pipe = pipe.to("cpu")
    if len(title.split(' ')) > 77:
        prompt = ' '.join(title.split(' ')[:77])
    image = pipe(title).images[0]
    image_path = f"{title.replace(' ', '_')}_cover.png"
    image.save(image_path)
    return image_path
def generate_presenter_notes(content):
    prompt = f"As a speaker, how would you explain this content:\n\n{content}\n\nSpeaker:"
    response = openai.Completion.create(
        engine="text-davinci-003",
        prompt=prompt,
        temperature=0.3,
        max_tokens=120
    )
    notes = response.choices[0].text.strip()
    notes = re.sub('[^a-zA-Z0-9 \n\.]', '', notes) # Remove special characters
    global total_tokens  # Access the global variable
    total_tokens += len(tokenizer.encode(notes))  # Update total tokens with the count of tokens in the API response
    return notes

def select_image(content, images):
    doc = nlp(content)
    words = [token.text for token in doc if token.pos_ == 'NOUN']
    most_common_noun = Counter(words).most_common(1)[0][0]
    for image, image_text in images.items():
        if most_common_noun in image_text:
            return image
    return None

def add_image_to_slide(slide, image_data_or_path):
    if isinstance(image_data_or_path, bytes):
        img_stream = BytesIO(image_data_or_path)
    else:  # Assuming it's a filepath if it's not bytes
        with open(image_data_or_path, 'rb') as f:
            image_data = f.read()
        img_stream = BytesIO(image_data)
    slide.shapes.add_picture(img_stream, Inches(2), Inches(2), height=Inches(5))
def add_bullet_points_to_slide(slide, points):
    bullet_slide = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(6), Inches(6)).text_frame
    for point in points:
        p = bullet_slide.add_paragraph()
        p.text = point
        p.level = 0

def create_presentation(titles, contents, images, image_indices, cover_image_path):
    presentation = Presentation()
    slide_layout = presentation.slide_layouts[0]
    slide = presentation.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "Your Presentation Title"
    # Adding the cover image
    add_image_to_slide(presentation.slides[0], cover_image_path)

    # Verify the length of titles and contents
    assert len(titles) == len(contents), "Mismatch between titles and contents length"

    for i in range(len(titles)):  # Iterate over the length of the titles/contents
        slide_layout = presentation.slide_layouts[1]
        slide = presentation.slides.add_slide(slide_layout)
        title = slide.shapes.title
        title.text = titles[i]
        add_bullet_points_to_slide(slide, contents[i])
        if i in image_indices:
            add_image_to_slide(slide, images[image_indices.index(i)])
        notes_slide = slide.notes_slide
        text_frame = notes_slide.notes_text_frame
        text_frame.text = generate_presenter_notes(contents[i])
    presentation.save("my_presentation.pptx")
    return presentation

def main():
    file_path = "document.pdf"
    pages = extract_text_from_pdf(file_path)  # Extract text from the pdf
    images, image_indices = extract_images_from_pdf(file_path)

    titles = []
    contents = []

    for page in pages:
        cleaned_page = preprocess_text(page)
        summary = generate_summary(cleaned_page)
        title = generate_title(cleaned_page, summary)
        titles.append(title)
        contents.append(summary)

    cover_image_path = generate_cover(titles[0])  # Generate cover using the title of the first slide

    presentation = create_presentation(titles, contents, images, image_indices, cover_image_path)

    print(f"Task completed in {round(time.time() - start_time, 2)/60} minutes")

if __name__ == "__main__":
    main()
